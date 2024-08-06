import os
import json
from pypdf import PdfReader
from pptx import Presentation
import google.generativeai as genai
from tkinter import *
from tkinter import filedialog

api_key = "Your API Key"
genai.configure(api_key=api_key)
model_gen = genai.GenerativeModel('gemini-1.5-flash')

def select_file():
    root = Tk()
    root.withdraw()  

    file_path = filedialog.askopenfilename(
        title="Pick a PDF or PPTX file",
        filetypes=[
            ("PDF and PPTX files", "*.pdf;*.pptx"),
            ("All files", "*.*")
        ]
    )

    root.destroy()  
    return file_path

file_path = select_file()
if file_path:
    print("File selected:", file_path)
else:
    print("No file chosen.")

input_file = file_path

number_of_cards = input("How many cards do you want? ")

def extract_text_from_pdf(file_path):
    reader = PdfReader(file_path)
    text = []
    for page in reader.pages:
        text.append(page.extract_text())
    return " ".join(text)

def extract_text_from_pptx(file_path):
    presentation = Presentation(file_path)
    text = []
    for slide in presentation.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

def extract_text(file_path):
    if file_path.lower().endswith('.pdf'):
        return extract_text_from_pdf(file_path)
    elif file_path.lower().endswith('.pptx'):
        return extract_text_from_pptx(file_path)
    else:
        print("Nope, this file format ain't supported.")
        return ""

def generate_flashcards(text):
    response = model_gen.generate_content(
        f"Make {number_of_cards} flashcards. Reply with just a JSON array. Each item should have 'front' and 'back'. Here's the text: {text}"
    )    

    if hasattr(response, 'text'):
        result = response.text
    else:
        result = response

    def clean_json_response(response):
        start_marker = "```json"
        end_marker = "```"

        if start_marker in response:
            response = response.split(start_marker)[-1]

        if end_marker in response:
            response = response.split(end_marker)[0]

        return response.strip()

    cleaned_result = clean_json_response(result)

    try:
        flashcards = json.loads(cleaned_result)
    except json.JSONDecodeError as e:
        print(f"Oops, JSON error: {e}")
        flashcards = []

    return flashcards

def save_flashcards_to_txt(flashcards, output_file):
    with open(output_file, 'w', encoding='utf-8') as file:
        for card in flashcards:
            file.write(f"{card['front']}:{card['back']}\n")
    print(f"Saved flashcards to {output_file}")

text = extract_text(input_file)

if text:
    print("Got the text! Sending it to AI.")
    flashcards = generate_flashcards(text)

    print("Flashcards done!")

    file_name = input('Name your .txt file (no extension): ')

    downloads_folder = os.path.expanduser("~/Downloads")
    output_file = os.path.join(downloads_folder, f'{file_name}.txt')

    save_flashcards_to_txt(flashcards, output_file)
else:
    print("No text found in the file.")

print ("If you like this, buy me a coffee: https://buymeacoffee.com/engineerdom!")
